from pptx import Presentation
import json

def pptx_to_json(pptx_path):
    prs = Presentation(pptx_path)
    ppt_data = {"slides": []}

    for slide_num, slide in enumerate(prs.slides, start=1):
        slide_data = {"slide_number": slide_num, "elements": []}

        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_data["elements"].append({
                    "type": "text",
                    "content": shape.text.strip(),
                    "position": {"left": shape.left, "top": shape.top}
                })

        ppt_data["slides"].append(slide_data)

    return ppt_data

ppt_json = pptx_to_json("./public/test01.pptx")
with open("./result/json/output01.json", "w", encoding="utf-8") as f:
    json.dump(ppt_json, f, indent=4, ensure_ascii=False)
